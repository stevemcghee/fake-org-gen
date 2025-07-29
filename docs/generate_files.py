

import argparse
import os
import random
import json
import io
import re
from faker import Faker
from docx import Document
import xlsxwriter
from pptx import Presentation
from PIL import Image
from google import genai
from google.genai import types

def sanitize_filename(title):
    """Converts a title into a safe filename."""
    # Remove invalid characters
    sanitized = re.sub(r'[\\/*?:"<>|]', "", title)
    # Replace spaces with underscores
    sanitized = sanitized.replace(" ", "_")
    # Truncate to a reasonable length
    return sanitized[:100]

def generate_unique_gemini_content(theme, file_type):
    """Generates unique, themed content for a specific file type by calling the Gemini API."""
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY environment variable not set.")
    
    client = genai.Client(api_key=api_key)
    
    if file_type == "document":
        doc_type = random.choice(["Internal Memo", "Project Proposal", "Competitive Analysis", "Budget Report", "Meeting Minutes"])
        content_prompt = f'- "doc_title": "A title for a {doc_type}".\n- "doc_body": "A 6-paragraph body for the document, relevant to the theme and the document type of {doc_type}".'
    elif file_type == "spreadsheet":
        sheet_type = random.choice(["Financial Statement", "Project Timeline", "Sales Tracker", "Inventory List", "Employee Directory"])
        content_prompt = f'- "sheet_title": "A title for a {sheet_type}".\n- "sheet_headers": "A list of 4-6 relevant column headers for the {sheet_type}".\n- "sheet_data": "A list of 15 lists, where each inner list is a row of realistic data for the {sheet_type}".'
    elif file_type == "presentation":
        ppt_type = random.choice(["Quarterly Review", "New Product Pitch", "Market Trend Analysis", "Team Training Guide"])
        content_prompt = f'- "ppt_title": "A title for a {ppt_type}".\n- "ppt_slide_details": "A JSON object with 4 slide titles as keys and a JSON list of 3-5 short, concise bullet points for each slide\u0027s body, relevant to a {ppt_type}".'
    elif file_type == "image":
        content_prompt = '- "image_prompt": "A concise, highly creative, and descriptive prompt for an AI image model to generate a photorealistic and thematic image, ensuring the prompt is unique and not a repeat of previous requests.".'
    else:
        return None

    prompt = f"""
    You are a creative assistant that generates realistic, unique business documents for a fictional company.
    The business theme is: "{theme}"
    The requested file type is: "{file_type}"

    Generate content for the following items in valid JSON format:
    {content_prompt}
    """
    try:
        response = client.models.generate_content(
            model="gemini-1.5-flash",
            contents=prompt
        )
        clean_response = response.text.strip().replace("```json", "").replace("```", "")
        return json.loads(clean_response)
    except (json.JSONDecodeError, Exception) as e:
        print(f"Error parsing Gemini text response: {e}")
        return None

def generate_image_from_api(prompt, file_path):
    """Generates an image using the Gemini API."""
    try:
        api_key = os.getenv("GEMINI_API_KEY")
        client = genai.Client(api_key=api_key)
        print(f"Submitting unique image prompt to Gemini: '{prompt}'")
        response = client.models.generate_content(
            model="gemini-2.0-flash-preview-image-generation",
            contents=prompt,
            config=types.GenerateContentConfig(response_modalities=['TEXT', 'IMAGE'])
        )
        for part in response.candidates[0].content.parts:
            if part.inline_data:
                image = Image.open(io.BytesIO(part.inline_data.data))
                image.save(file_path)
                print(f"Successfully generated and saved AI image: {file_path}")
                return
    except Exception as e:
        print(f"Could not generate or save image due to an API error: {e}")
        img = Image.new('RGB', (800, 600), color=(50, 10, 10))
        from PIL import ImageDraw
        d = ImageDraw.Draw(img)
        d.text((10, 10), f"Image generation failed.\nError: {e}", fill=(255, 255, 255))
        img.save(file_path)

def generate_document(user, org_name, file_path, fake, theme_content):
    doc = Document()
    doc.add_heading(f'{org_name} - {theme_content["doc_title"]}', 0)
    p = doc.add_paragraph(f'This document is the property of {org_name}. ')
    p.add_run('Prepared by: ').bold = True
    p.add_run(user)
    doc.add_paragraph()
    doc_body = theme_content.get("doc_body", "")
    if isinstance(doc_body, list):
        for paragraph in doc_body:
            doc.add_paragraph(paragraph)
    elif isinstance(doc_body, str):
        for paragraph in doc_body.split('\n'):
            doc.add_paragraph(paragraph)
    doc.save(file_path)

def generate_spreadsheet(user, org_name, file_path, fake, theme_content):
    wb = xlsxwriter.Workbook(file_path)
    ws = wb.add_worksheet()
    bold = wb.add_format({'bold': True})
    ws.write('A1', f'{org_name} - {theme_content["sheet_title"]}', bold)
    ws.write('A2', 'Prepared by:', bold)
    ws.write('B2', user)
    headers = theme_content.get("sheet_headers", [])
    for col, header in enumerate(headers):
        ws.write(4, col, header, bold)
    sheet_data = theme_content.get("sheet_data", [])
    for row_num, row_data in enumerate(sheet_data, 5):
        for col_num, cell_data in enumerate(row_data):
            ws.write(row_num, col_num, cell_data)
    wb.close()

def generate_presentation(user, org_name, file_path, fake, theme_content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f'{org_name} - {theme_content["ppt_title"]}'
    subtitle.text = f"Presented by {user}\n{fake.date()}"
    
    bullet_slide_layout = prs.slide_layouts[1]
    for slide_title, bullet_points in theme_content.get("ppt_slide_details", {}).items():
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = slide_title
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear existing text
        
        if isinstance(bullet_points, list):
            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
        else: # Fallback for unexpected format
            p = tf.add_paragraph()
            p.text = str(bullet_points)

    prs.save(file_path)

def main():
    parser = argparse.ArgumentParser(description="Generate unique, random files for a small business using the Gemini API.")
    parser.add_argument("--users", nargs="+", default=["frodobaggins", "gandalf"], help="List of usernames.")
    parser.add_argument("--num-files", type=int, default=5, help="Number of files per user.")
    parser.add_argument("--org-name", default="Shire Holdings", help="Name of the organization.")
    parser.add_argument("--theme", required=True, help="Business model theme.")
    args = parser.parse_args()

    fake = Faker()

    file_generators = {
        "document": generate_document,
        "spreadsheet": generate_spreadsheet,
        "presentation": generate_presentation,
        "image": lambda user, org, path, f, content: generate_image_from_api(content.get('image_prompt'), path),
    }
    file_extensions = {
        "document": ".docx", "spreadsheet": ".xlsx",
        "presentation": ".pptx", "image": ".png",
    }

    for user in args.users:
        output_dir = os.path.join("output", user)
        os.makedirs(output_dir, exist_ok=True)
        for i in range(args.num_files):
            file_type = random.choice(list(file_generators.keys()))
            
            print(f"--- Generating unique content for {file_type} ---")
            theme_content = generate_unique_gemini_content(args.theme, file_type)
            
            if theme_content:
                title = theme_content.get("doc_title") or theme_content.get("sheet_title") or theme_content.get("ppt_title") or f"image_{i+1}"
                file_name = f"{sanitize_filename(title)}{file_extensions[file_type]}"
                file_path = os.path.join(output_dir, file_name)
                
                print(f"--- Saving file: {file_path} ---")
                try:
                    file_generators[file_type](user, args.org_name, file_path, fake, theme_content)
                except Exception as e:
                    print(f"Error generating file {file_path}: {e}")
            else:
                print(f"Skipping file due to content generation failure.")

if __name__ == "__main__":
    main()
